import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from src.state import AgentState
from src.utils.logger import logger

def hex_to_rgb(hex_color: str):
    """Converts a hex color string to RGBColor object."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return RGBColor(0, 0, 0) # default to black
    return RGBColor(*tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4)))

def ppt_builder_agent(state: AgentState) -> dict:
    """
    Agent 10: Uses python-pptx to generate slides, applying styles, content, and diagrams.
    """
    logger.info("Agent 10: PPT Builder Agent started.")
    content = state.get('slide_content', [])
    visual = state.get('visual_design_guidelines', {})
    diagrams = state.get('diagrams', {})
    
    colors = visual.get('color_palette', ['#FFFFFF', '#000000', '#AAAAAA', '#FF0000'])
    primary = hex_to_rgb(colors[0])
    
    prs = Presentation()
    
    for slide_data in content:
        # Choose layout (1: title slide, 5: simple text, etc.)
        # Here we just use a blank layout or a title + body layout
        layout = prs.slide_layouts[1] # Title and Body
        slide = prs.slides.add_slide(layout)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data.get('slide_title', 'Untitled')
        if title.text_frame:
            for p in title.text_frame.paragraphs:
                p.font.color.rgb = primary
                # p.font.name = visual.get('font_heading')
        
        # Body
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear() # clear existing
        
        points = slide_data.get('body_points', [])
        for pt in points:
            p = tf.add_paragraph()
            p.text = str(pt)
            p.font.size = Pt(24)
            p.level = 0
            
        # Notes
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame
        notes_text.text = slide_data.get('speaker_notes', '')
        
        # Check if architecture slide to insert diagram
        if 'architecture' in str(title.text).lower() and 'architecture' in diagrams:
            try:
                # Add picture
                pic_path = diagrams['architecture']
                if os.path.exists(pic_path):
                    slide.shapes.add_picture(pic_path, Inches(1), Inches(2.5), width=Inches(8))
            except Exception as e:
                logger.error(f"Failed to add image to PPT: {str(e)}")
            
    import time
    
    # Save Presentation
    os.makedirs('output', exist_ok=True)
    
    timestamp = time.strftime("%Y%m%d-%H%M%S")
    ppt_filename = f"presentation_{timestamp}.pptx"
    ppt_path = os.path.abspath(f'output/{ppt_filename}')
    
    try:
        prs.save(ppt_path)
    except PermissionError:
        # Fallback if somehow even that name has an issue
        ppt_path = os.path.abspath(f'output/presentation_fallback_{time.time()}.pptx')
        prs.save(ppt_path)
        
    logger.info(f"Agent 10 saved presentation to {ppt_path}")
    
    return {"pptx_path": ppt_path, "current_agent": "Agent 10"}
