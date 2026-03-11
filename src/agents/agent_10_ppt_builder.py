import os
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from src.state import AgentState
from src.utils.logger import logger

def hex_to_rgb(hex_code):
    """Convert HEX string to RGBColor object."""
    hex_code = hex_code.lstrip('#')
    if len(hex_code) == 6:
        return RGBColor(int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16))
    return RGBColor(0, 0, 0) # Default to black if invalid

def ppt_builder_agent(state: AgentState) -> dict:
    """
    Agent 10: Uses python-pptx to generate slides, applying styles, content, and diagrams.
    """
    logger.info("Agent 10: PPT Builder Agent started.")
    content = state.get('slide_content', [])
    theme = state.get('visual_design_guidelines', {})
    diagrams = state.get('diagrams', {})
    topic = state.get('topic', 'ShrutiBuilder Pitch')
    
    if not content:
        return {"errors": state.get('errors', []) + ["Agent 10: No slide content found."]}

    # Extract Theme
    bg_hex = theme.get('background_color', '#FFFFFF')
    primary_hex = theme.get('primary_color', '#0F52BA')
    font_name = theme.get('font_family', 'Calibri')
    
    # Determine text color based on background (simple bright/dark check would be ideal, but we'll try to infer or just set bright text on dark pg)
    # Defaulting to white text for dark modes, dark text for light modes
    is_dark_mode = True if bg_hex.lower() in ['#000000', '#111111', '#1a1a1a', '#222222', '#121212', '#0f0f0f'] else False
    if bg_hex.startswith('#0') or bg_hex.startswith('#1') or bg_hex.startswith('#2'):
        is_dark_mode = True
        
    text_color = RGBColor(255, 255, 255) if is_dark_mode else RGBColor(50, 50, 50)
    bg_rgb = hex_to_rgb(bg_hex)
    primary_rgb = hex_to_rgb(primary_hex)

    prs = Presentation()
    
    # Optional: adjust slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for idx, slide_data in enumerate(content):
        # Choose layout (0: title slide, 1: title + content, etc.)
        layout_index = 0 if idx == 0 else 1
        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)
        
        # Manually apply Background Color via a huge rectangle behind everything
        background = slide.shapes.add_shape(
            1, # msoShapeRectangle
            0, 0, prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = bg_rgb
        background.line.fill.background() # No outline
        slide.shapes._spTree.remove(background._element)
        slide.shapes._spTree.insert(2, background._element) # Send to back
        
        # Add decorative bar at top for style
        header_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.15))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = primary_rgb
        header_bar.line.fill.background()
        
        # Title
        if slide.shapes.title:
            title = slide.shapes.title
            title.text = slide_data.get('slide_title', 'Untitled')
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.name = font_name
                paragraph.font.color.rgb = primary_rgb
                paragraph.font.bold = True
                if idx == 0:
                    paragraph.font.size = Pt(54)
                else:
                    paragraph.font.size = Pt(40)
                    paragraph.alignment = PP_ALIGN.LEFT
                    
            if idx != 0:
                title.left = Inches(0.5)
                title.top = Inches(0.5)
                title.width = Inches(12)
                title.height = Inches(1)
            
        # Body
        if len(slide.shapes.placeholders) > 1:
            body = slide.shapes.placeholders[1]
            if not body.has_text_frame:
                continue
                
            if idx != 0:
                body.left = Inches(0.5)
                body.top = Inches(1.8)
                body.width = Inches(12)
                body.height = Inches(5)
                
            text_frame = body.text_frame
            points = slide_data.get('body_points', [])
            
            # Format text explicitly
            text_frame.text = "" # Clear placeholder
            for pt in points:
                p = text_frame.add_paragraph()
                p.text = pt
                p.font.name = font_name
                p.font.color.rgb = text_color
                p.font.size = Pt(24)
                p.space_after = Pt(14) # Spacing
                
        # Notes
        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame
            notes_text.text = slide_data.get('speaker_notes', '')
        except Exception:
            pass # Non-critical if notes slide fails
            
        # Check for matching diagrams
        if f"slide_{idx}" in diagrams:
            pic_path = diagrams[f"slide_{idx}"]
            if os.path.exists(pic_path):
                try:
                    # Place image at right side usually
                    slide.shapes.add_picture(pic_path, Inches(7), Inches(1.8), height=Inches(4.5))
                except Exception as e:
                    logger.error(f"Failed to add image to PPT: {str(e)}")
                    
    # Save Presentation
    timestamp = int(time.time())
    
    os.makedirs('output', exist_ok=True)
    ppt_filename = f"presentation_{timestamp}.pptx"
    ppt_path = os.path.abspath(f'output/{ppt_filename}')
    
    try:
        prs.save(ppt_path)
        logger.info(f"Agent 10 successfully wrote stylish PPTX to: {ppt_path}")
        return {"pptx_path": ppt_path, "current_agent": "Agent 10 PPT"}
    except Exception as e:
        logger.error(f"Agent 10 failed to save PPTX: {str(e)}")
        # Try a fallback if permission error
        ppt_path = os.path.abspath(f'output/presentation_fallback_{time.time()}.pptx')
        prs.save(ppt_path)
        return {"pptx_path": ppt_path, "current_agent": "Agent 10 PPT"}
